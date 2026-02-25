"""
Agent 工具系统

工具（Tool）是 Agent 与外界交互的桥梁。
Agent 通过调用工具来获取信息或执行操作。

工具系统的核心概念：
1. BaseTool: 工具基类，定义工具的基本结构
2. ToolRegistry: 工具注册中心，管理所有可用工具
3. 工具发现: Agent 如何知道有哪些工具可用

工具调用流程：
┌──────────┐    选择工具    ┌──────────┐    执行    ┌──────────┐
│  Agent   │ ────────────→ │  Tool    │ ────────→ │  Result  │
└──────────┘               └──────────┘           └──────────┘
     ↑                                                   │
     └───────────────── 返回结果 ────────────────────────┘

示例：
用户: "帮我算一下 2+2"
Agent: 思考... 用户需要计算，我有 calculator 工具
Agent: 调用 calculator(expression="2+2")
工具: 返回 "计算结果: 4"
Agent: 根据结果回答 "2+2 = 4"
"""

import json
from typing import Callable, Dict, List, Optional, Any
from abc import ABC, abstractmethod
from pydantic import BaseModel, Field
import logging

logger = logging.getLogger(__name__)


class ToolSchema(BaseModel):
    """
    工具参数 Schema

    定义工具输入参数的结构，用于：
    1. 验证输入参数是否正确
    2. 生成工具描述给 LLM（让 LLM 知道如何调用）

    示例：
    class CalculatorSchema(ToolSchema):
        expression: str = Field(description="要计算的数学表达式")
    """
    pass


class BaseTool(ABC):
    """
    工具基类

    所有 Agent 工具都需要继承这个基类。

    创建新工具的步骤：
    1. 继承 BaseTool
    2. 设置 name（工具名称）和 description（工具描述）
    3. 定义 args_schema（参数结构）
    4. 实现 _run 方法（工具的实际执行逻辑）

    示例：
    class CalculatorTool(BaseTool):
        name = "calculator"
        description = "计算数学表达式"

        class ArgsSchema(ToolSchema):
            expression: str = Field(description="数学表达式，如 '2+2'")

        args_schema = ArgsSchema

        def _run(self, expression: str) -> str:
            result = eval(expression)
            return f"计算结果: {result}"
    """

    # 工具名称（唯一标识）
    # 用于 Agent 识别和调用工具
    # 示例: "calculator", "web_search"
    name: str = ""

    # 工具描述
    # 这个描述会被发送给 LLM，帮助 LLM 理解工具的用途
    # 描述要清晰、具体，让 LLM 能正确选择工具
    # 示例: "计算数学表达式，如加法、减法、乘法、除法"
    description: str = ""

    # 参数 Schema
    # 定义工具需要的输入参数及其类型
    # 使用 Pydantic 定义，会自动生成 JSON Schema
    args_schema: Optional[type[ToolSchema]] = None

    @abstractmethod
    def _run(self, **kwargs) -> str:
        """
        执行工具的核心逻辑

        子类必须实现这个方法。

        Args:
            **kwargs: 工具输入参数，由 args_schema 定义结构

        Returns:
            工具执行结果（字符串格式）
        """
        pass

    def run(self, **kwargs) -> str:
        """
        执行工具（带错误处理）

        这是工具的入口方法，会调用 _run 并处理异常。

        Args:
            **kwargs: 工具输入参数

        Returns:
            工具执行结果或错误信息
        """
        try:
            # 验证参数（如果有 schema）
            if self.args_schema:
                validated = self.args_schema(**kwargs)
                kwargs = validated.model_dump()

            # 执行工具
            result = self._run(**kwargs)
            logger.info(f"工具 {self.name} 执行成功: {result[:100]}...")
            return result

        except Exception as e:
            error_msg = f"工具 {self.name} 执行失败: {str(e)}"
            logger.error(error_msg)
            return error_msg

    def to_openai_tool(self) -> Dict[str, Any]:
        """
        转换为 OpenAI Tool 格式

        OpenAI 的 Function Calling 使用特定的 JSON 格式来描述工具。
        这个方法将工具定义转换为 OpenAI 兼容的格式。

        Returns:
            OpenAI Tool 格式的字典

        示例输出：
        {
            "type": "function",
            "function": {
                "name": "calculator",
                "description": "计算数学表达式",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "expression": {
                            "type": "string",
                            "description": "数学表达式"
                        }
                    },
                    "required": ["expression"]
                }
            }
        }
        """
        tool_def = {
            "type": "function",
            "function": {
                "name": self.name,
                "description": self.description,
            }
        }

        # 添加参数 schema
        if self.args_schema:
            # 使用 Pydantic 生成 JSON Schema
            schema = self.args_schema.model_json_schema()
            # 移除不需要的字段
            schema.pop("title", None)
            tool_def["function"]["parameters"] = schema

        return tool_def

    def get_tool_prompt(self) -> str:
        """
        获取工具的文本描述

        用于在 Prompt 中告诉 LLM 有哪些工具可用。

        Returns:
            工具的文本描述
        """
        prompt = f"- {self.name}: {self.description}\n"
        if self.args_schema:
            schema = self.args_schema.model_json_schema()
            props = schema.get("properties", {})
            if props:
                prompt += "  参数:\n"
                for prop_name, prop_def in props.items():
                    prop_desc = prop_def.get("description", "")
                    prompt += f"    - {prop_name}: {prop_desc}\n"
        return prompt


class ToolRegistry:
    """
    工具注册中心

    管理所有 Agent 可用的工具。

    职责：
    1. 注册工具
    2. 获取工具
    3. 列出所有工具
    4. 生成工具描述（给 LLM）

    使用示例：
    registry = ToolRegistry()
    registry.register(CalculatorTool())

    # 获取工具
    tool = registry.get_tool("calculator")
    result = tool.run(expression="2+2")

    # 获取所有工具的 OpenAI 格式
    tools = registry.get_openai_tools()
    """

    def __init__(self):
        """初始化工具注册中心"""
        # 工具存储：{工具名称: 工具实例}
        self._tools: Dict[str, BaseTool] = {}

    def register(self, tool: BaseTool) -> None:
        """
        注册工具

        Args:
            tool: 工具实例

        Raises:
            ValueError: 如果工具名称已存在
        """
        if tool.name in self._tools:
            raise ValueError(f"工具 '{tool.name}' 已存在")

        self._tools[tool.name] = tool
        logger.info(f"已注册工具: {tool.name}")

    def get_tool(self, name: str) -> Optional[BaseTool]:
        """
        获取工具

        Args:
            name: 工具名称

        Returns:
            工具实例，如果不存在返回 None
        """
        return self._tools.get(name)

    def list_tools(self) -> List[str]:
        """
        列出所有已注册的工具名称

        Returns:
            工具名称列表
        """
        return list(self._tools.keys())

    def get_openai_tools(self) -> List[Dict[str, Any]]:
        """
        获取所有工具的 OpenAI 格式

        用于传递给 LLM 进行 Function Calling。

        Returns:
            OpenAI Tool 格式的列表
        """
        return [tool.to_openai_tool() for tool in self._tools.values()]

    def get_tools_prompt(self) -> str:
        """
        获取所有工具的文本描述

        用于在 Prompt 中告诉 LLM 有哪些工具可用。

        Returns:
            所有工具的文本描述
        """
        prompts = ["可用工具列表：\n"]
        for tool in self._tools.values():
            prompts.append(tool.get_tool_prompt())
        return "\n".join(prompts)

    def execute_tool(self, name: str, arguments: Dict[str, Any]) -> str:
        """
        执行工具

        根据工具名称和参数执行工具，返回结果。

        Args:
            name: 工具名称
            arguments: 工具参数

        Returns:
            工具执行结果

        Raises:
            ValueError: 如果工具不存在
        """
        tool = self.get_tool(name)
        if not tool:
            raise ValueError(f"工具 '{name}' 不存在")

        return tool.run(**arguments)


# ==================== 内置工具 ====================

class CalculatorTool(BaseTool):
    """
    计算器工具

    用于计算数学表达式，如加减乘除、幂运算等。

    这是一个简单的演示工具，用于验证 Agent 框架是否正常工作。
    """

    name = "calculator"
    description = "计算数学表达式。支持加减乘除、幂运算等基本数学运算。例如：'2+2'、'10*5'、'2**10'"

    class ArgsSchema(ToolSchema):
        expression: str = Field(
            description="要计算的数学表达式，如 '2+2'、'10*5'、'sqrt(16)' 等"
        )

    args_schema = ArgsSchema

    def _run(self, expression: str) -> str:
        """
        执行计算

        注意：这里使用 eval() 仅用于演示目的。
        在生产环境中，应该使用更安全的表达式解析库。

        Args:
            expression: 数学表达式

        Returns:
            计算结果
        """
        # 安全检查：只允许数字、运算符和常用数学函数
        import math
        allowed_names = {
            "abs": abs, "round": round, "min": min, "max": max,
            "sqrt": math.sqrt, "pow": math.pow, "sin": math.sin,
            "cos": math.cos, "tan": math.tan, "log": math.log,
            "log10": math.log10, "exp": math.exp, "pi": math.pi, "e": math.e,
        }

        # 简单的安全检查
        expression = expression.strip()
        if not all(c.isdigit() or c in "+-*/.() **%sqrtabscoundminmaxrtpwelg10xp" for c in expression.lower()):
            return "错误：表达式包含不允许的字符"

        try:
            # 使用受限的命名空间执行
            result = eval(expression, {"__builtins__": {}}, allowed_names)
            return f"计算结果: {result}"
        except Exception as e:
            return f"计算错误: {str(e)}"


class EchoTool(BaseTool):
    """
    回显工具

    简单地将输入返回，用于测试 Agent 是否能正确调用工具。
    """

    name = "echo"
    description = "回显用户输入的内容。用于测试工具调用是否正常。"

    class ArgsSchema(ToolSchema):
        message: str = Field(description="要回显的消息")

    args_schema = ArgsSchema

    def _run(self, message: str) -> str:
        return f"Echo: {message}"


# ==================== Skill 工具适配器 ====================

class SkillToolAdapter(BaseTool):
    """
    技能工具适配器

    将 Skill 包装成 Tool，让 Agent 可以通过工具调用使用技能。

    这个适配器的作用：
    1. 将技能注册为 Agent 可调用的工具
    2. 同步调用技能的异步执行方法
    3. 转换技能参数为工具参数格式

    示例：
        skill = CalculatorSkill()
        tool = SkillToolAdapter(skill)
        registry.register(tool)
        # 现在 Agent 可以通过工具调用技能了
    """

    def __init__(self, skill):
        """
        初始化技能工具适配器

        Args:
            skill: BaseSkill 实例
        """
        self._skill = skill
        self.name = skill.name
        self.description = skill.description

        # 从技能配置构建参数 schema
        self.args_schema = self._build_args_schema(skill)

    def _build_args_schema(self, skill):
        """
        从技能配置构建参数 schema

        使用 Pydantic v2 的 create_model 动态创建模型。

        Args:
            skill: BaseSkill 实例

        Returns:
            Pydantic Model 类
        """
        from pydantic import create_model

        parameters = skill.config.parameters

        if not parameters:
            return None

        # 使用 pydantic.create_model 动态创建模型
        # Pydantic v2 格式: field_name: (type, Field(...))
        fields = {}
        for param_name, param_config in parameters.items():
            param_type = str  # 默认字符串类型
            param_desc = param_config.get("description", "")
            param_required = param_config.get("required", False)

            if param_required:
                # 必填字段：使用 Field(...) 作为默认值
                fields[param_name] = (
                    param_type, Field(description=param_desc))
            else:
                # 可选字段：使用 Field(default="") 作为默认值
                fields[param_name] = (param_type, Field(
                    default="", description=param_desc))

        # 动态创建 Pydantic 模型
        return create_model(
            f"{skill.name.capitalize()}ArgsSchema",
            __base__=ToolSchema,
            **fields
        )

    def _run(self, **kwargs) -> str:
        """
        执行技能（同步包装）

        由于 Agent 工具是同步的，我们需要在事件循环中执行异步技能。

        Args:
            **kwargs: 技能参数

        Returns:
            技能执行结果
        """
        import asyncio

        try:
            # 尝试在现有事件循环中运行
            loop = asyncio.get_event_loop()
            if loop.is_running():
                # 如果事件循环正在运行，创建一个新的任务
                import concurrent.futures
                with concurrent.futures.ThreadPoolExecutor() as executor:
                    future = executor.submit(
                        asyncio.run,
                        self._skill.execute(**kwargs)
                    )
                    return future.result()
            else:
                # 如果没有运行的事件循环，直接运行
                return loop.run_until_complete(self._skill.execute(**kwargs))
        except RuntimeError:
            # 没有事件循环，创建一个新的
            return asyncio.run(self._skill.execute(**kwargs))


def register_skills_as_tools(skill_registry, tool_registry):
    """
    将所有技能注册为工具

    这个函数将技能注册中心的技能包装成工具，注册到工具注册中心。
    这样 Agent 就可以通过工具调用使用技能了。

    Args:
        skill_registry: 技能注册中心
        tool_registry: 工具注册中心

    Returns:
        注册的工具数量
    """
    count = 0
    for skill in skill_registry.list_skills(enabled_only=True):
        try:
            # 创建工具适配器
            tool = SkillToolAdapter(skill)

            # 检查工具是否已存在
            if tool.name not in tool_registry.list_tools():
                tool_registry.register(tool)
                count += 1
                # 打印工具描述，便于调试
                tool_prompt = tool.get_tool_prompt()
                logger.info(f"已将技能 {skill.name} 注册为工具")
                logger.debug(f"工具描述:\n{tool_prompt}")
            else:
                logger.debug(f"工具 {tool.name} 已存在，跳过注册")

        except Exception as e:
            logger.error(f"注册技能 {skill.name} 为工具失败: {e}")

    # 打印所有已注册的工具
    logger.info(f"当前已注册的工具: {tool_registry.list_tools()}")

    return count


# ==================== 全局工具注册中心 ====================

# 创建全局工具注册中心实例
# 在应用启动时注册所有内置工具
global_tool_registry = ToolRegistry()


# ==================== 文件读取工具 ====================

class FileReadTool(BaseTool):
    """
    文件读取工具

    读取本地文件内容，支持多种文件格式。
    Agent 使用此工具读取用户上传的附件文件。
    """

    name = "file_read"
    description = "读取本地文件内容。用于分析用户上传的附件文件。支持文本文件、PDF、Markdown、代码文件等格式。"

    class ArgsSchema(ToolSchema):
        file_path: str = Field(description="文件的完整路径")
        max_length: int = Field(default=10000, description="最大读取字符数，默认10000")

    args_schema = ArgsSchema

    def _run(self, file_path: str, max_length: int = 10000) -> str:
        """
        读取文件内容

        Args:
            file_path: 文件路径
            max_length: 最大读取字符数

        Returns:
            文件内容或错误信息
        """
        import os

        try:
            # 检查文件是否存在
            if not os.path.exists(file_path):
                return f"错误：文件不存在 - {file_path}"

            # 获取文件扩展名
            _, ext = os.path.splitext(file_path)
            ext = ext.lower()

            # 获取文件大小
            file_size = os.path.getsize(file_path)

            # 根据文件类型选择读取方式
            if ext == ".pdf":
                return self._read_pdf(file_path, max_length)
            elif ext in [".doc", ".docx"]:
                return self._read_docx(file_path, max_length)
            elif ext in [".ppt", ".pptx"]:
                return self._read_pptx(file_path, max_length)
            elif ext in [".xls", ".xlsx"]:
                return self._read_xlsx(file_path, max_length)
            else:
                # 默认作为文本文件读取
                return self._read_text(file_path, max_length, file_size)

        except Exception as e:
            return f"读取文件失败: {str(e)}"

    def _read_text(self, file_path: str, max_length: int, file_size: int) -> str:
        """读取文本文件"""
        import os

        try:
            # 尝试不同编码
            encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
            content = None

            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read(max_length + 1)
                    break
                except UnicodeDecodeError:
                    continue

            if content is None:
                return f"错误：无法解码文件内容，可能不是文本文件"

            # 检查是否截断
            truncated = "" if len(
                content) <= max_length else "\n\n[...文件内容已截断...]"

            return f"文件大小: {file_size} 字节\n\n{content[:max_length]}{truncated}"

        except Exception as e:
            return f"读取文本文件失败: {str(e)}"

    def _read_pdf(self, file_path: str, max_length: int) -> str:
        """读取PDF文件"""
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(file_path)
            total_pages = len(doc)  # 先保存总页数
            content_parts = []
            total_length = 0

            for page_num, page in enumerate(doc):
                if total_length >= max_length:
                    content_parts.append(f"\n[...已截断，共{total_pages}页...]")
                    break

                text = page.get_text()
                content_parts.append(f"=== 第{page_num + 1}页 ===\n{text}")
                total_length += len(text)

            doc.close()

            full_content = "\n\n".join(content_parts)
            return f"PDF文件，共{total_pages}页\n\n{full_content[:max_length]}"

        except ImportError:
            return "错误：PDF读取需要安装 PyMuPDF 库 (pip install pymupdf)"
        except Exception as e:
            return f"读取PDF失败: {str(e)}"

    def _read_docx(self, file_path: str, max_length: int) -> str:
        """读取Word文档"""
        try:
            from docx import Document

            doc = Document(file_path)
            content_parts = []
            total_length = 0

            for para in doc.paragraphs:
                if total_length >= max_length:
                    content_parts.append("\n[...内容已截断...]")
                    break
                content_parts.append(para.text)
                total_length += len(para.text)

            return f"Word文档\n\n" + "\n".join(content_parts)[:max_length]

        except ImportError:
            return "错误：Word文档读取需要安装 python-docx 库 (pip install python-docx)"
        except Exception as e:
            return f"读取Word文档失败: {str(e)}"

    def _read_pptx(self, file_path: str, max_length: int) -> str:
        """读取PPT文件"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            content_parts = []
            total_length = 0

            for slide_num, slide in enumerate(prs.slides):
                if total_length >= max_length:
                    content_parts.append("\n[...内容已截断...]")
                    break

                slide_content = [f"=== 幻灯片 {slide_num + 1} ==="]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_content.append(shape.text)

                text = "\n".join(slide_content)
                content_parts.append(text)
                total_length += len(text)

            return f"PPT文件，共{len(prs.slides)}页\n\n" + "\n\n".join(content_parts)[:max_length]

        except ImportError:
            return "错误：PPT读取需要安装 python-pptx 库 (pip install python-pptx)"
        except Exception as e:
            return f"读取PPT失败: {str(e)}"

    def _read_xlsx(self, file_path: str, max_length: int) -> str:
        """读取Excel文件"""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(file_path, read_only=True)
            content_parts = []
            total_length = 0

            for sheet_name in wb.sheetnames:
                if total_length >= max_length:
                    content_parts.append("\n[...内容已截断...]")
                    break

                sheet = wb[sheet_name]
                sheet_content = [f"=== 工作表: {sheet_name} ==="]

                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(
                        str(cell) if cell is not None else "" for cell in row)
                    sheet_content.append(row_text)
                    total_length += len(row_text)

                    if total_length >= max_length:
                        break

                content_parts.append("\n".join(sheet_content))

            wb.close()

            return f"Excel文件，共{len(wb.sheetnames)}个工作表\n\n" + "\n\n".join(content_parts)[:max_length]

        except ImportError:
            return "错误：Excel读取需要安装 openpyxl 库 (pip install openpyxl)"
        except Exception as e:
            return f"读取Excel失败: {str(e)}"


def init_default_tools():
    """
    初始化默认工具

    在应用启动时调用，注册所有内置工具。
    """
    global global_tool_registry

    # 注册内置工具
    global_tool_registry.register(CalculatorTool())
    global_tool_registry.register(EchoTool())
    global_tool_registry.register(FileReadTool())  # 文件读取工具

    logger.info(f"已注册 {len(global_tool_registry.list_tools())} 个默认工具")


# 应用启动时初始化工具
init_default_tools()
